#!/usr/bin/env python3
"""Build the EE4615 Group 22 presentation (2-D Vernier TDC) on the TU Delft template."""
import copy
import os
from PIL import Image, ImageChops
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = "/home/danieltyukov/workspace/tud/tud-digital-ic-design-2"
TEMPLATE = ("/home/danieltyukov/workspace/tud/tud-masters-thesis-pick/"
            "TU Delft - Powerpoint templates/TU Delft - Corporate Presentation-ENG_v1.5.pptx")
OUT = os.path.join(ROOT, "slides", "EE4615_Group22_2D_Vernier_TDC.pptx")
PREP = os.path.join(ROOT, "slides", "figs", "prepared")
os.makedirs(PREP, exist_ok=True)

NAVY = RGBColor(0x0C, 0x23, 0x40)
CYAN = RGBColor(0x00, 0xA6, 0xD6)
GRAY = RGBColor(0x5B, 0x67, 0x70)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xE9, 0xF5, 0xFB)
ROWALT = RGBColor(0xF2, 0xF5, 0xF7)
GREEN = RGBColor(0x00, 0x8A, 0x4B)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
ML = Inches(0.79)          # left margin
CW = Inches(11.76)         # content width


# ---------------------------------------------------------------- figure prep
def trim_white(src, dst, pad=12, bg=255):
    im = Image.open(src).convert("RGB")
    diff = ImageChops.difference(im, Image.new("RGB", im.size, (bg, bg, bg)))
    bbox = diff.getbbox()
    if bbox:
        l, t, r, b = bbox
        l = max(0, l - pad); t = max(0, t - pad)
        r = min(im.width, r + pad); b = min(im.height, b + pad)
        im = im.crop((l, t, r, b))
    im.save(dst)
    return dst


def trim_dark(src, dst, pad=16, thresh=40):
    im = Image.open(src).convert("RGB")
    gray = im.convert("L").point(lambda p: 255 if p > thresh else 0)
    bbox = gray.getbbox()
    if bbox:
        l, t, r, b = bbox
        l = max(0, l - pad); t = max(0, t - pad)
        r = min(im.width, r + pad); b = min(im.height, b + pad)
        im = im.crop((l, t, r, b))
    im.save(dst)
    return dst


FIG = {}
for name, rel in {
    "overview": "slides/figs/overview_svg.png",
    "srlatch": "slides/figs/srlatch_svg.png",
    "tau1": "slides/figs/delay_tau1_svg.png",
    "tau2": "slides/figs/delay_tau2_svg.png",
    "staircase": "results/corners/staircase_all_corners.png",
    "dnl": "results/corners/dnl_per_corner.png",
    "inl": "results/corners/inl_per_corner.png",
    "lsb_tuned": "results/tuned_corners/lsb_tuned_vs_untuned.png",
    "calrange": "results/tuned_corners/calibration_range.png",
    "meta": "results/metastability/metastability_curve.png",
    "deadzone": "results/metastability/deadzone_transfer.png",
    "scaling": "results/vernier1d_vs_2d/scaling_1d_vs_2d.png",
    "dnl1d2d": "results/vernier1d_vs_2d/dnl_1d_vs_2d.png",
}.items():
    FIG[name] = trim_white(os.path.join(ROOT, rel), os.path.join(PREP, name + ".png"))

for name, rel in {
    "core": "slides/figs/tdc_core_canvas.png",
    "core_zoom": "slides/figs/tdc_core_zoom_canvas.png",
    "tb": "slides/figs/tb_tdc_canvas.png",
}.items():
    FIG[name] = trim_dark(os.path.join(ROOT, rel), os.path.join(PREP, name + ".png"))


# ---------------------------------------------------------------- helpers
prs = Presentation(TEMPLATE)
# drop the template's sample slides (rels too, so the parts are not serialized)
for sid in list(prs.slides._sldIdLst):
    prs.part.drop_rel(sid.get(qn('r:id')))
    prs.slides._sldIdLst.remove(sid)

LAYOUTS = {l.name: l for l in prs.slide_layouts}
L_TITLE = LAYOUTS["Title slide + image"]
L_PLAIN = LAYOUTS["Alleen titel"]
L_CLOSE = LAYOUTS["Closure"]

page_no = [0]


def set_run(r, text, size, bold=False, color=NAVY, italic=False, name="Arial"):
    r.text = text
    r.font.name = name
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return r


def textbox(slide, x, y, w, h, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def hang(p, left_in, hang_in):
    pPr = p._p.get_or_add_pPr()
    pPr.set('marL', str(int(Inches(left_in))))
    pPr.set('indent', str(-int(Inches(hang_in))))


def bullets(slide, x, y, w, h, items, size=16.0, sub_size=14.0, space=8):
    """items: list of (level, segments). segments = str or list of (text, bold, color)."""
    tb, tf = textbox(slide, x, y, w, h)
    first = True
    for level, segs in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(space)
        if level == 0:
            hang(p, 0.26, 0.26)
            set_run(p.add_run(), "▪  ", size, bold=True, color=CYAN)
            sz = size
        else:
            hang(p, 0.56, 0.20)
            set_run(p.add_run(), "–  ", sub_size, bold=True, color=GRAY)
            sz = sub_size
        if isinstance(segs, str):
            segs = [(segs, False, NAVY)]
        for text, bold, color in segs:
            set_run(p.add_run(), text, sz, bold=bold, color=color)
    return tb


def pic_caption(slide, pic, text):
    y = pic.top + pic.height + Inches(0.06)
    return caption(slide, pic.left, y, pic.width, text)


def accent_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, ML, Inches(1.45), Inches(2.6), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CYAN
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def presenter_tag(slide, who):
    tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.45), Inches(0.74),
                                 Inches(1.10), Inches(0.34))
    tag.adjustments[0] = 0.5
    tag.fill.solid()
    tag.fill.fore_color.rgb = CYAN
    tag.line.fill.background()
    tag.shadow.inherit = False
    tf = tag.text_frame
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(), who, 11, bold=True, color=WHITE)


def content_slide(title, presenter=None):
    slide = prs.slides.add_slide(L_PLAIN)
    ti = slide.shapes.title
    ti.left, ti.top, ti.width, ti.height = ML, Inches(0.72), Inches(10.5), Inches(0.5)
    tf = ti.text_frame
    p = tf.paragraphs[0]
    set_run(p.add_run(), title, 25, bold=True, color=NAVY)
    accent_bar(slide)
    page_no[0] += 1
    tb, tfn = textbox(slide, Inches(12.35), Inches(7.02), Inches(0.7), Inches(0.3))
    pn = tfn.paragraphs[0]
    pn.alignment = PP_ALIGN.RIGHT
    set_run(pn.add_run(), str(page_no[0]), 11, color=GRAY)
    if presenter:
        presenter_tag(slide, presenter)
    return slide


def add_picture_fit(slide, path, x, y, w, h, border=False):
    im = Image.open(path)
    ar = im.width / im.height
    bw, bh = w, h
    if bw / bh > ar:
        nw, nh = int(bh * ar), bh
    else:
        nw, nh = bw, int(bw / ar)
    nx = int(x + (bw - nw) / 2)
    ny = int(y + (bh - nh) / 2)
    pic = slide.shapes.add_picture(path, nx, ny, nw, nh)
    if border:
        pic.line.color.rgb = RGBColor(0xC9, 0xD2, 0xD8)
        pic.line.width = Pt(0.75)
    return pic


def caption(slide, x, y, w, text, align=PP_ALIGN.CENTER):
    tb, tf = textbox(slide, x, y, w, Inches(0.3))
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), text, 10.5, italic=True, color=GRAY)
    return tb


def stat_card(slide, x, y, w, h, big, small, big_size=24):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = CARD
    box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.06)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(), big, big_size, bold=True, color=CYAN)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(2)
    set_run(p2.add_run(), small, 11, color=NAVY)
    return box


def make_table(slide, x, y, w, rows, col_widths, header, data,
               size=12.5, header_size=12.5, row_h=0.32, highlight_col=None,
               col_aligns=None):
    shape = slide.shapes.add_table(len(data) + 1, len(header), x, y, w, Inches(row_h * (len(data) + 1)))
    table = shape.table
    # disable banding, use our own fills
    tbl = shape._element.graphic.graphicData.tbl
    tbl[0].set('firstRow', '0')
    tbl[0].set('bandRow', '0')
    for i, cw in enumerate(col_widths):
        table.columns[i].width = cw
    for r in range(len(data) + 1):
        table.rows[r].height = Inches(row_h)
    for c, htxt in enumerate(header):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
        p = cell.text_frame.paragraphs[0]
        if col_aligns:
            p.alignment = col_aligns[c]
        set_run(p.add_run(), htxt, header_size, bold=True, color=WHITE)
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 == 0 else ROWALT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            p = cell.text_frame.paragraphs[0]
            if col_aligns:
                p.alignment = col_aligns[c]
            bold = False
            color = NAVY
            if isinstance(val, tuple):
                val, bold, color = val
            if highlight_col is not None and c == highlight_col:
                bold = True
            set_run(p.add_run(), val, size, bold=bold, color=color)
    return shape


# ================================================================ slide 1: title
slide = prs.slides.add_slide(L_TITLE)
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 10:
        tf = ph.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        set_run(p.add_run(), "A 2-D Vernier Time-to-Digital Converter", 31, bold=True, color=WHITE)
    elif idx == 11:
        tf = ph.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        set_run(p.add_run(), "15 ps resolution, 5-bit thermometer output, TSMC 180 nm BCD",
                15, bold=True, color=NAVY)
        p2 = tf.add_paragraph(); p2.space_before = Pt(8)
        set_run(p2.add_run(), "EE4615 Digital IC Design II  |  Group 22", 12.5, color=NAVY)
        p3 = tf.add_paragraph(); p3.space_before = Pt(2)
        set_run(p3.add_run(), "Daniel Tyukov (5714699)  |  Joris Hoogeweegen (5619963)", 12.5, color=NAVY)
        p4 = tf.add_paragraph(); p4.space_before = Pt(2)
        set_run(p4.add_run(), "Delft University of Technology  |  10 June 2026", 12.5, color=NAVY)
    elif idx == 13:
        ph.insert_picture(FIG["core"])

# ================================================================ slide 2: the brief
slide = content_slide("The brief: quantize a time interval to 5 bits", "Daniel")
bullets(slide, ML, Inches(1.75), CW, Inches(3.3), [
    (0, [("Task: ", True, NAVY),
         ("convert the delay between a START and a STOP edge into a 5-bit thermometer code Q₁..Q₃₁", False, NAVY)]),
    (0, [("Hard requirements: ", True, NAVY),
         ("resolution better than 20 ps, full 32-code range, and validated robustness in all five process corners", False, NAVY)]),
    (0, [("Rules of the game: ", True, NAVY),
         ("TSMC 180 nm BCD at 1.8 V, every gate hand-built (no standard cells), schematic plus Spectre only", False, NAVY)]),
    (0, [("Optimize the trade-off: ", True, NAVY),
         ("linearity (DNL, INL), energy per conversion, area as total transistor width ΣW, and FoM = energy per step", False, NAVY)]),
    (0, [("Our choice: ", True, NAVY),
         ("a two-dimensional Vernier TDC (Vercesi et al., JSSC 2010), the structure that decouples resolution from range", False, NAVY)]),
])
cy = Inches(5.55)
cwd = Inches(2.76); gap = Inches(0.24)
specs = [("5 bit", "thermometer output"), ("< 20 ps", "required resolution"),
         ("5 / 5", "corners must pass"), ("0", "standard cells allowed")]
for i, (big, small) in enumerate(specs):
    stat_card(slide, ML + i * (cwd + gap), cy, cwd, Inches(1.05), big, small)

# ================================================================ slide 3: why 2-D vernier
slide = content_slide("Why Vernier, and why a 2-D grid", "Joris")
bullets(slide, ML, Inches(1.75), Inches(6.1), Inches(4.9), [
    (0, [("Vernier principle: ", True, NAVY),
         ("START runs through slow stages τ₁, STOP chases through fast stages τ₂", False, NAVY)]),
    (1, "the LSB is the difference t₀ = τ₁ − τ₂, far below a single gate delay"),
    (0, [("1-D catch: ", True, NAVY),
         ("range = N·t₀, so 32 codes need 32 stages in each line; cells and latency grow linearly", False, NAVY)]),
    (0, [("2-D idea: ", True, NAVY),
         ("compare START tap i against STOP tap j; every pair measures Δt(i,j) = i·τ₁ − j·τ₂", False, NAVY)]),
    (1, "a grid of N+N stages yields ~N²/2 distinct codes: range without long lines"),
    (0, [("Payoff for our 5-bit target: ", True, NAVY),
         ("31 codes from only 10 + 6 delay stages, one SR-latch arbiter per code", False, NAVY)]),
], size=15.5, sub_size=13.5)
make_table(slide, Inches(7.25), Inches(1.95), Inches(5.3),
           4, [Inches(2.10), Inches(1.60), Inches(1.60)],
           ["For 32 codes", "1-D Vernier", "2-D (ours)"],
           [["Delay stages", "32 + 32", ("10 + 6", True, CYAN)],
            ["Arbiters", "32", "31 + 29 dummy"],
            ["Full-scale latency", "~3.3 ns", ("~0.9 ns", True, CYAN)],
            ["LSB", "τ₁ − τ₂", "τ₁ − τ₂"]],
           size=13, header_size=13, row_h=0.42)
caption(slide, Inches(7.25), Inches(4.15), Inches(5.3),
        "Same resolution, a quarter of the delay cells")

# ================================================================ slide 4: architecture
slide = content_slide("Architecture: a 10 × 6 grid locked at k = 6", "Daniel")
bullets(slide, ML, Inches(1.75), Inches(6.7), Inches(4.9), [
    (0, [("Operating point: ", True, NAVY),
         ("τ₁ = 6·t₀ = 90 ps (10 stages), τ₂ = 5·t₀ = 75 ps (6 stages), so t₀ = 15 ps with 5 ps of spec margin", False, NAVY)]),
    (0, [("Bijective routing: ", True, NAVY),
         ("each thermometer level m maps to exactly one grid cell, so no OR-tree and no skew from combining logic", False, NAVY)]),
    (1, "yₘ = ((m−1) mod 6) + 1,   xₘ = (m + 5·yₘ) / 6"),
    (0, [("Matched loading: ", True, NAVY),
         ("31 active SR latches plus 29 dummies pad the array to 10 × 10, every tap drives identical fan-out", False, NAVY)]),
    (0, [("I/O: ", True, NAVY),
         ("buffer trees distribute START, STOP and the 40 ps RESET; the code is read straight off the latch outputs", False, NAVY)]),
], size=15.5, sub_size=13.5)
add_picture_fit(slide, FIG["overview"], Inches(7.7), Inches(1.6), Inches(4.9), Inches(4.95))
caption(slide, Inches(7.7), Inches(6.65), Inches(4.9),
        "Two delay lines, SR-latch matrix, dummy padding")

# ================================================================ slide 5: delay stage
slide = content_slide("Delay stage: FO = 1 inverter pair with MOSCAP trim", "Joris")
bullets(slide, ML, Inches(1.70), CW, Inches(2.4), [
    (0, [("Breaking the FO = 4 habit: ", True, NAVY),
         ("classical sizing pinned the stage floor at 97 ps; an FO = 1 pair opened a 74 to 147 ps tuning window", False, NAVY)]),
    (0, [("Twin cells: ", True, NAVY),
         ("both lines use identical active devices; the 6:5 delay ratio comes from loading, so τ₁ and τ₂ track across PVT", False, NAVY)]),
    (0, [("Calibration built in: ", True, NAVY),
         ("a fixed MOSCAP sets the nominal delay, three transmission-gate switched banks re-center t₀ per corner (authority ≈ 0 to 48 ps)", False, NAVY)]),
], size=15, sub_size=13.5, space=6)
add_picture_fit(slide, FIG["tau1"], Inches(1.6), Inches(3.85), Inches(7.4), Inches(3.05))
add_picture_fit(slide, FIG["tau2"], Inches(9.25), Inches(3.85), Inches(3.4), Inches(3.05))
caption(slide, Inches(1.6), Inches(6.95), Inches(7.4),
        "delay_tau1 (START, 90 ps): trim<0:2> banks behind transmission gates", align=PP_ALIGN.CENTER)
caption(slide, Inches(9.25), Inches(6.95), Inches(3.4), "delay_tau2 (STOP, 75 ps)")

# ================================================================ slide 6: arbiter
slide = content_slide("Arbiter: cross-coupled NAND SR latch, 13 transistors", "Daniel")
bullets(slide, ML, Inches(1.70), CW, Inches(2.4), [
    (0, [("Topology: ", True, NAVY),
         ("two 2-input NANDs (all devices 440 nm / 180 nm), an async-reset NMOS, and buffered Q / Q̅ outputs", False, NAVY)]),
    (0, [("Race logic: ", True, NAVY),
         ("both NAND outputs idle high; the first rising input wins, so the latch directly records which edge arrived first", False, NAVY)]),
    (0, [("Why it works: ", True, NAVY),
         ("equal N and P widths balance the S and R paths, keeping the arbitration offset at the femtosecond level", False, NAVY)]),
    (0, [("Why not a D flip-flop: ", True, NAVY),
         ("60 instances make cell economy count; the SR latch is smaller, symmetric, and resets in one transistor", False, NAVY)]),
], size=15, sub_size=13.5, space=6)
pic = add_picture_fit(slide, FIG["srlatch"], Inches(2.0), Inches(3.20), Inches(9.3), Inches(3.55))
pic_caption(slide, pic, "vernier2d/srlatch: S and R race, RESET clears the grid in 40 ps")

# ================================================================ slide 7: loading pivot
slide = content_slide("The pivot: size against the real load, not the ideal one", "Joris")
bullets(slide, ML, Inches(1.75), Inches(6.6), Inches(4.9), [
    (0, [("The trap: ", True, NAVY),
         ("a grid tap drives a full row or column of latches plus the next delay stage, about a 20× equivalent load", False, NAVY)]),
    (0, [("We measured it: ", True, NAVY),
         ("a 1-D row built from the same cells with raw tap loading predicted t₀ = 16.6 ps, but measured 28.6 ps", False, NAVY)]),
    (1, "sizing against an idealized load gives a fictitious LSB"),
    (0, [("The fix: ", True, NAVY),
         ("characterize every stage under the replica 20× load, and equalize all taps with dummy latches", False, NAVY)]),
    (0, [("Consequence: ", True, NAVY),
         ("the original k = 4 point (60 / 45 ps) was unreachable; the design re-locked at k = 6 (90 / 75 ps)", False, NAVY)]),
    (0, [("Lesson: ", True, NAVY),
         ("in a Vernier grid, loading is part of the architecture, not a layout afterthought", False, NAVY)]),
], size=15, sub_size=13.5)
pic = add_picture_fit(slide, FIG["core_zoom"], Inches(7.45), Inches(2.45), Inches(5.3), Inches(3.0))
pic_caption(slide, pic, "Grid detail: DUMMY latches pad every tap to fan-out 10")

# ================================================================ slide 8: built + testbench
slide = content_slide("Built in Cadence, verified with the course testbench", "Daniel")
bullets(slide, ML, Inches(1.70), CW, Inches(2.15), [
    (0, [("The core: ", True, NAVY),
         ("11 + 7 delay stages, 60 SR latches, 1120 devices, ΣW = 1306 µm, all cells hand-built in vernier2d", False, NAVY)]),
    (0, [("Honest energy numbers: ", True, NAVY),
         ("the DUT sits on its own supply behind a series ammeter; stimulus buffers run from a separate rail", False, NAVY)]),
    (0, [("Automated sweeps: ", True, NAVY),
         ("OCEAN steps Δt through all 32 codes per corner and extracts LSB, DNL, INL, energy and pass / fail", False, NAVY)]),
], size=15, sub_size=13.5, space=6)
pic = add_picture_fit(slide, FIG["core"], Inches(1.45), Inches(3.35), Inches(10.4), Inches(3.35))
pic_caption(slide, pic, "vernier2d/tdc_core: τ₁ line (top), τ₂ line (left), 10 × 10 padded latch matrix")

# ================================================================ slide 9: staircase
slide = content_slide("Transfer: a clean 5-bit staircase in every corner", "Joris")
add_picture_fit(slide, FIG["staircase"], ML, Inches(1.8), Inches(6.4), Inches(4.7))
make_table(slide, Inches(7.55), Inches(1.95), Inches(4.9),
           6, [Inches(1.30), Inches(1.80), Inches(1.80)],
           ["Corner", "LSB (ps)", "Spec < 20 ps"],
           [["TT", "15.35", ("PASS", True, GREEN)],
            ["SS", "18.20", ("PASS", True, GREEN)],
            ["FF", "13.10", ("PASS", True, GREEN)],
            ["SF", "15.35", ("PASS", True, GREEN)],
            ["FS", "15.70", ("PASS", True, GREEN)]],
           size=13, header_size=13, row_h=0.40)
bullets(slide, Inches(7.55), Inches(4.75), Inches(5.0), Inches(2.0), [
    (0, "One fixed trim code, no per-corner tweaking"),
    (0, "All 31 steps present in all five corners"),
    (0, "Offset stays a static, calibratable shift"),
], size=14.5, space=6)
caption(slide, ML, Inches(6.55), Inches(6.4), "Measured TDC transfer, all corners, 27 °C, 1.8 V")

# ================================================================ slide 10: linearity
slide = content_slide("Linearity: no missing codes anywhere", "Daniel")
add_picture_fit(slide, FIG["dnl"], ML, Inches(1.85), Inches(5.75), Inches(4.1))
add_picture_fit(slide, FIG["inl"], Inches(6.85), Inches(1.85), Inches(5.75), Inches(4.1))
caption(slide, ML, Inches(6.0), Inches(5.75), "DNL per corner (LSB)")
caption(slide, Inches(6.85), Inches(6.0), Inches(5.75), "INL per corner (LSB)")
bullets(slide, ML, Inches(6.35), CW, Inches(0.9), [
    (0, [("Worst DNL 0.90 LSB (SS), far from the −1 LSB missing-code limit; worst INL 1.04 LSB, all other corners below 0.8", False, NAVY)]),
    (0, [("The every-6-codes DNL signature is the row-to-column residual of the grid, the expected 2-D Vernier fingerprint", False, NAVY)]),
], size=14, space=4)

# ================================================================ slide 11: calibration
slide = content_slide("Per-corner MOSCAP calibration: 6× tighter LSB", "Joris")
add_picture_fit(slide, FIG["lsb_tuned"], ML, Inches(1.85), Inches(5.75), Inches(4.1))
add_picture_fit(slide, FIG["calrange"], Inches(6.85), Inches(1.85), Inches(5.75), Inches(4.1))
caption(slide, ML, Inches(6.0), Inches(5.75), "LSB per corner, untuned vs calibrated")
caption(slide, Inches(6.85), Inches(6.0), Inches(5.75), "Trim-bank authority: t₀ vs configuration")
bullets(slide, ML, Inches(6.35), CW, Inches(0.9), [
    (0, [("Three trim banks per stage act as static configuration bits, set once per corner", False, NAVY)]),
    (0, [("LSB corner spread shrinks from 5.10 ps to 0.85 ps; every corner re-centers to 15.0 – 15.9 ps, worst INL drops to 0.74 LSB", False, NAVY)]),
], size=14, space=4)

# ================================================================ slide 12: metastability
slide = content_slide("How small a Δt can the arbiter resolve?", "Joris")
add_picture_fit(slide, FIG["meta"], ML, Inches(1.85), Inches(5.75), Inches(4.1))
add_picture_fit(slide, FIG["deadzone"], Inches(6.85), Inches(1.85), Inches(5.75), Inches(4.1))
caption(slide, ML, Inches(6.0), Inches(5.75), "Decision time vs overdrive: τ = 24.1 ps regeneration")
caption(slide, Inches(6.85), Inches(6.0), Inches(5.75), "Arbiter transfer: dead zone below 1 fs")
bullets(slide, ML, Inches(6.35), CW, Inches(0.9), [
    (0, [("Arbitration offset +0.15 ps (1 % of one LSB); decision ambiguity below 1 fs, four orders under the 15 ps LSB", False, NAVY)]),
    (0, [("Worst observed decision time 348 ps at 0.5 fs overdrive: metastability does not limit this design, the LSB does", False, NAVY)]),
], size=14, space=4)

# ================================================================ slide 13: 1D vs 2D
slide = content_slide("Measured: 2-D beats 1-D built from the same cells", "Daniel")
add_picture_fit(slide, FIG["scaling"], ML, Inches(1.85), Inches(5.9), Inches(4.35))
make_table(slide, Inches(7.05), Inches(1.95), Inches(5.45),
           6, [Inches(2.35), Inches(1.55), Inches(1.55)],
           ["Same cells, measured", "1-D row", "2-D core"],
           [["LSB (ps)", "28.57", ("15.35", True, CYAN)],
            ["DNL max (LSB)", "0.89", ("0.66", True, CYAN)],
            ["ΣW per code (µm)", "81.4", ("42.1", True, CYAN)],
            ["Energy per code (pJ)", "0.327", "0.376"],
            ["Latency, 31 codes", "3.3 ns", ("0.9 ns", True, CYAN)]],
           size=13, header_size=13, row_h=0.40)
bullets(slide, Inches(7.05), Inches(4.65), Inches(5.5), Inches(1.9), [
    (0, "Raw tap loading costs the 1-D row almost 2× in resolution"),
    (0, "Scaling to 31 codes: 1-D needs 3.4× the stages and 1.9× the area"),
], size=14, space=6)
caption(slide, ML, Inches(6.45), Inches(5.9), "Cost scaling with code count, 1-D vs 2-D")

# ================================================================ slide 14: scorecard
slide = content_slide("Scorecard: every requirement met, with margin", "Joris")
make_table(slide, ML, Inches(1.80), Inches(11.76),
           10, [Inches(3.60), Inches(2.70), Inches(2.90), Inches(2.56)],
           ["Requirement", "Target", "Achieved", "Verdict"],
           [["Output bits", "5", "5 (thermometer Q₁..Q₃₁)", ("PASS", True, GREEN)],
            ["Resolution (LSB)", "< 20 ps", "15.0 – 15.9 ps calibrated", ("PASS", True, GREEN)],
            ["Corner robustness", "TT, SS, FF, SF, FS", "5 / 5 corners pass", ("PASS", True, GREEN)],
            ["DNL", "> −1 LSB (no missing codes)", "worst 0.89 LSB", ("PASS", True, GREEN)],
            ["INL", "report", "worst 0.74 LSB (calibrated)", ("PASS", True, GREEN)],
            ["Energy per conversion", "minimize", "10.7 – 13.0 pJ", ("✓", True, GREEN)],
            ["FoM (energy per step)", "minimize", "0.48 – 0.53 pJ / step", ("✓", True, GREEN)],
            ["Area ΣW", "minimize", "1305.6 µm (1120 devices)", ("✓", True, GREEN)],
            ["ENOB", "report", "4.40 – 4.62 bit", ("✓", True, GREEN)]],
           size=12.5, header_size=12.5, row_h=0.42,
           col_aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.CENTER])
caption(slide, ML, Inches(6.05), Inches(11.76),
        "All numbers from Spectre corner sweeps at 27 °C, 1.8 V, course TB_TDC harness", align=PP_ALIGN.LEFT)

# ================================================================ slide 15: conclusions
slide = content_slide("Conclusions and outlook", "Daniel")
tb, tf = textbox(slide, ML, Inches(1.70), Inches(6.3), Inches(0.4))
set_run(tf.paragraphs[0].add_run(), "What we showed", 15, bold=True, color=CYAN)
bullets(slide, ML, Inches(2.15), Inches(6.3), Inches(2.6), [
    (0, "A 15 ps, 5-bit 2-D Vernier TDC in 180 nm that meets every course spec in all five corners"),
    (0, "Two ideas carried the design: size under the real grid load, and equalize fan-out with dummies"),
    (0, "Calibration is cheap: three MOSCAP banks per stage buy a 6× tighter LSB spread"),
], size=14.5, space=6)
tb, tf = textbox(slide, ML, Inches(4.55), Inches(6.3), Inches(0.4))
set_run(tf.paragraphs[0].add_run(), "Future work", 15, bold=True, color=CYAN)
bullets(slide, ML, Inches(5.00), Inches(6.3), Inches(1.7), [
    (0, "DLL-based background calibration to replace static trim codes"),
    (0, "Voltage-to-time front end: reuse the TDC as an ADC quantizer"),
    (0, "Layout and extracted re-simulation"),
], size=14.5, space=6)
box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.55), Inches(1.70), Inches(5.0), Inches(3.55))
box.fill.solid(); box.fill.fore_color.rgb = CARD
box.line.fill.background(); box.shadow.inherit = False
tb, tf = textbox(slide, Inches(7.85), Inches(1.95), Inches(4.4), Inches(3.2))
p = tf.paragraphs[0]
set_run(p.add_run(), "The design in six numbers", 15, bold=True, color=CYAN)
entries = [
    ("15 ps", "resolution t₀, spec asks < 20 ps"),
    ("0.85 ps", "LSB corner spread after calibration"),
    ("11.7 pJ", "energy per conversion (TT)"),
    ("0.49 pJ / step", "figure of merit (TT)"),
    ("1306 µm", "total transistor width ΣW"),
    ("< 1 fs", "arbiter dead zone"),
]
for num, what in entries:
    p = tf.add_paragraph(); p.space_before = Pt(9)
    set_run(p.add_run(), num + "   ", 14.5, bold=True, color=CYAN)
    set_run(p.add_run(), what, 13.5, color=NAVY)
tb, tf = textbox(slide, Inches(7.55), Inches(5.55), Inches(5.0), Inches(0.8))
set_run(tf.paragraphs[0].add_run(),
        "All design, simulation and analysis were carried out jointly by both team members.",
        12, italic=True, color=GRAY)

# ================================================================ slide 16: closure
slide = prs.slides.add_slide(L_CLOSE)
for ph in list(slide.placeholders):
    if ph.has_text_frame and ph.placeholder_format.idx != 0:
        continue
tb, tf = textbox(slide, ML, Inches(2.6), Inches(8.0), Inches(1.0))
set_run(tf.paragraphs[0].add_run(), "Thank you. Questions?", 30, bold=True, color=WHITE)
tb, tf = textbox(slide, ML, Inches(3.7), Inches(8.0), Inches(1.2))
p = tf.paragraphs[0]
set_run(p.add_run(), "Group 22: Daniel Tyukov, Joris Hoogeweegen", 14, color=WHITE)
p2 = tf.add_paragraph(); p2.space_before = Pt(4)
set_run(p2.add_run(), "2-D Vernier TDC, EE4615 Digital IC Design II", 14, color=WHITE)

prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides.slide_renderable_part_list) if hasattr(prs.slides,'slide_renderable_part_list') else len(prs.slides._sldIdLst))
